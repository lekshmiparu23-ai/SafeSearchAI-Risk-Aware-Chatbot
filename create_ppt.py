import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("python-pptx is not installed. Please run: pip install python-pptx")
    sys.exit(1)

prs = Presentation()

def add_slide(title_text, content_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title_text
    
    tf = slide.placeholders[1].text_frame
    for i, bullet in enumerate(content_bullets):
        if i == 0:
            tf.text = bullet
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Risk-Aware Query Classification for Secure AI Chatbot Systems"
subtitle.text = "Lekshmi Maniyan, Mohammed Shameer M, Mrs. P.V. Indira\n\nDhanalakshmi Srinivasan University"

# Slide 2: The Problem
add_slide("The Problem with Current AI Safety", [
    "Existing AI safety systems face significant challenges:",
    "• Limitation 1: Static keyword filters are too rigid, permanently blocking safe queries (False Positives).",
    "• Limitation 2: Standard LLMs suffer from 'artificial confidence inflation', confidently delivering harmful instructions.",
    "• Danger: Leads to successful jailbreaks, prompt injections, and illegal instructional generation."
])

# Slide 3: Our Solution
add_slide("Our Solution: SafeSearchAI", [
    "A novel pipeline to secure large language models:",
    "• Introduces a 5-layer hierarchical safety architecture.",
    "• Implements independent safety guarantees at each processing step.",
    "• Core Innovation: Constrained Risk-Aware Scoring.",
    "• Guarantees mathematically bounded monotonic alignment between risk and confidence metrics."
])

# Slide 4: System Architecture
add_slide("System Architecture", [
    "A resilient edge-to-cloud implementation:",
    "• Native Android Application (Kotlin / Jetpack Compose) for mobile inference.",
    "• Lightning-speed Llama-3.1-8B LLM processing powered by Cerebras Inference API.",
    "• Real-time cloud-native data persistence via Firebase Firestore.",
    "\n[Please Insert Architecture Diagram Figure 1 from Paper Here]"
])

# Slide 5: Layers 1 & 2
add_slide("Layers 1 & 2: The Frontline Defense", [
    "Layer 1: Deterministic Safety Override",
    "• Implements direct string heuristics to catch explicit threats instantly.",
    "• Contextual exception engine allows safe discussion of dual-use terms (e.g., medical 'bacteria').",
    "Layer 2: Semantic Intent Classification",
    "• Analyzes edge cases using Llama 3.1 to categorize into Safe, Sensitive, or Harmful categories."
])

# Slide 6: Layers 3, 4 & 5
add_slide("Layers 3, 4 & 5: Core Innovation", [
    "Layer 3: Policy Enforcement",
    "• Immediate query termination for classified 'Harmful' intents.",
    "Layer 4: Semantic Reliability Audit",
    "• Evaluates response on Relevance, Completeness, Specificity, Clarity, and Consistency (0-100 scale).",
    "Layer 5: Constrained Risk-Aware Scoring",
    "• S_final = min(S_rel, S_max(R_k))",
    "• Physically restricts the display confidence of 'Sensitive' queries, preventing user deception."
])

# Slide 7: Evaluation & Results
add_slide("Evaluation & Results", [
    "Rigorous testing against 1,050 complex and adversarial queries:",
    "• Achieved an overall Macro-averaged F1-Score of 94.1%.",
    "• Reached 98.1% Precision for Harmful queries.",
    "• Zero 'Harmful-to-Safe' critical misclassifications.",
    "• Significantly outperformed Google Perspective API and LlamaGuard baselines."
])

# Slide 8: Real-World Applications
add_slide("Real-World Applications", [
    "The 'Sensitive' intermediary tier allows unique deployments:",
    "• Enterprise Customer Support: gracefully handles off-topic or competitor inquiries.",
    "• Healthcare Information Systems: safely distinguishes between clinical advice and self-harm contexts.",
    "• Educational AI Assistants: allows academic exploration without providing dangerous operational guidelines."
])

# Slide 9: Conclusion & Future Work
add_slide("Conclusion & Future Work", [
    "Conclusion",
    "• SafeSearchAI provides the first mathematically bounded confidence scoring for conversational AI.",
    "Future Work",
    "• Multi-lingual LLM support.",
    "• Continued adversarial fine-tuning.",
    "• RAG-based factual grounding."
])

output_file = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SafeSearchAI_Presentation.pptx")
prs.save(output_file)
print(f"Presentation saved successfully to: {output_file}")
